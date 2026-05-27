"""
Extract plain text from PDFs, PPTX, DOCX, and Jupyter notebooks.

For each supported document, writes a sidecar text file:
    <name>.<ext>.extracted.txt   (default suffix)

Slides keep `--- Slide N: <title> ---` markers.
Notebooks keep `--- Cell N [type] ---` markers.

Idempotent: skips files whose sidecar exists and is newer than the source.

Usage:
    python tools/extract_docs.py --root . [--out-suffix .extracted.txt] [--force]
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

SUPPORTED = {".pdf", ".pptx", ".docx", ".ipynb"}


def needs_extract(src: Path, sidecar: Path, force: bool) -> bool:
    if force or not sidecar.exists():
        return True
    try:
        return sidecar.stat().st_mtime < src.stat().st_mtime
    except OSError:
        return True


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as e:
        raise RuntimeError(f"pypdf not installed: {e}") from e
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        try:
            txt = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            txt = f"<page extract error: {e}>"
        parts.append(f"--- Page {i} ---\n{txt.strip()}\n")
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError(f"python-pptx not installed: {e}") from e
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
        except Exception:  # noqa: BLE001
            pass
        body: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line and line != title:
                        body.append(line)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:  # noqa: BLE001
            pass
        head = f"--- Slide {i}: {title or '(untitled)'} ---"
        block = head + "\n" + "\n".join(body)
        if notes:
            block += f"\n\n[notes]\n{notes}"
        parts.append(block + "\n")
    return "\n".join(parts)


def extract_docx(path: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise RuntimeError(f"python-docx not installed: {e}") from e
    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts) + "\n"


def extract_ipynb(path: Path) -> str:
    try:
        nb = json.loads(path.read_text(encoding="utf-8"))
    except Exception as e:  # noqa: BLE001
        raise RuntimeError(f"could not parse notebook: {e}") from e
    parts: list[str] = []
    for i, cell in enumerate(nb.get("cells", []), 1):
        ctype = cell.get("cell_type", "?")
        src = cell.get("source", "")
        if isinstance(src, list):
            src = "".join(src)
        parts.append(f"--- Cell {i} [{ctype}] ---\n{src.rstrip()}\n")
    return "\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".ipynb": extract_ipynb,
}


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--root", default=".", help="Workspace root (default: cwd)")
    ap.add_argument("--out-suffix", default=".extracted.txt",
                    help="Suffix appended to source filename for the sidecar")
    ap.add_argument("--force", action="store_true", help="Re-extract even if sidecar is up-to-date")
    args = ap.parse_args(argv)

    root = Path(args.root).resolve()
    if not root.is_dir():
        print(f"error: root not found: {root}", file=sys.stderr)
        return 2

    targets: list[Path] = []
    for ext in SUPPORTED:
        targets.extend(root.rglob(f"*{ext}"))

    skip_dirs = {".git", "node_modules", ".venv", "venv", "__pycache__", "dist", "build", "out"}
    targets = [p for p in targets if not any(part in skip_dirs for part in p.parts)]

    done = skipped = errors = 0
    for src in sorted(targets):
        sidecar = src.with_name(src.name + args.out_suffix)
        if not needs_extract(src, sidecar, args.force):
            skipped += 1
            continue
        ext = src.suffix.lower()
        fn = EXTRACTORS.get(ext)
        if not fn:
            continue
        try:
            text = fn(src)
            sidecar.write_text(text, encoding="utf-8")
            print(f"OK   {src.relative_to(root)}  ->  {sidecar.name}")
            done += 1
        except Exception as e:  # noqa: BLE001
            print(f"FAIL {src.relative_to(root)} :: {e}", file=sys.stderr)
            errors += 1

    print(f"\nExtracted {done}, skipped {skipped}, failed {errors}.")
    return 1 if errors else 0


if __name__ == "__main__":
    raise SystemExit(main())
